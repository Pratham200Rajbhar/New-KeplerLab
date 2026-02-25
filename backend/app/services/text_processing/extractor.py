"""Unified text extraction from files, URLs, and YouTube.

Every extraction method returns a standard ``ExtractionResult`` dict:

    {
      "text":        str,
      "status":      "success" | "failed",
      "source":      str,
      "source_type": str,
      "title":       str | None,
      "metadata":    dict,
      "error":       str | None,   # only on failure
    }

New formats supported (added in this revision):
    EPUB, ODT, EML (email), MSG (Outlook), SVG, generic text fallback.
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any, Callable, Dict, Optional

import chardet

from .file_detector import FileTypeDetector

logger = logging.getLogger(__name__)

# Type alias
ExtractionResult = Dict[str, Any]

# ---------------------------------------------------------------------------
# Process-level model singletons (lazy-loaded, never re-created)
# ---------------------------------------------------------------------------
_ocr_service_instance: Optional[Any] = None
_transcription_service_instance: Optional[Any] = None


def _get_ocr_service() -> Any:
    global _ocr_service_instance
    if _ocr_service_instance is None:
        from .ocr_service import OCRService
        logger.info("Initialising OCRService singleton — first use.")
        _ocr_service_instance = OCRService()
    return _ocr_service_instance


def _get_transcription_service() -> Any:
    global _transcription_service_instance
    if _transcription_service_instance is None:
        from .transcription_service import AudioTranscriptionService
        logger.info("Initialising AudioTranscriptionService singleton — first use.")
        _transcription_service_instance = AudioTranscriptionService()
    return _transcription_service_instance


# ── Result builders ───────────────────────────────────────────────────────────


def _ok(
    text: str,
    source: str,
    source_type: str = "file",
    *,
    title: Optional[str] = None,
    **metadata,
) -> ExtractionResult:
    return {
        "text":        text,
        "status":      "success",
        "source":      source,
        "source_type": source_type,
        "title":       title,
        "word_count":  len(text.split()) if text else 0,
        "metadata":    metadata,
    }


def _fail(source: str, error: str, source_type: str = "file") -> ExtractionResult:
    return {
        "text":        "",
        "status":      "failed",
        "source":      source,
        "source_type": source_type,
        "error":       error,
        "metadata":    {},
    }


# ── Extractor registry ────────────────────────────────────────────────────────

_DOC_EXTRACTORS: Dict[str, Callable] = {}


def _dataframe_summary(df: "pd.DataFrame", label: str = "Dataset") -> str:
    """Build a RAG-friendly text representation of a DataFrame.

    Includes schema, descriptive statistics, and head/tail samples —
    enough for the LLM to reason about structure *and* content without
    needing the full data.
    """
    import pandas as pd

    lines: list[str] = [
        f"=== {label} ===",
        f"Shape: {df.shape[0]} rows × {df.shape[1]} columns",
        f"Columns: {', '.join(str(c) for c in df.columns)}",
        f"Column types: {', '.join(f'{c}: {t}' for c, t in df.dtypes.items())}",
        "",
        "Statistics:",
        df.describe(include="all").to_string(),
        "",
        "Sample (first 5 rows):",
        df.head(5).to_string(),
        "",
        "Sample (last 5 rows):",
        df.tail(5).to_string(),
    ]
    return "\n".join(lines)


def _register(*exts: str):
    """Decorator: register an extractor for one or more file extensions."""
    def wrapper(fn):
        for ext in exts:
            _DOC_EXTRACTORS[ext] = fn
        return fn
    return wrapper


# ── Format extractors ─────────────────────────────────────────────────────────


@_register("pdf")
def _extract_pdf(path: str) -> ExtractionResult:
    """PDF — PyMuPDF primary, PDFPlumber tables, OCR fallback for image-only pages."""
    try:
        from .pdf_extractor import PDFExtractor
        from .resilient_runner import run_with_retry
        from app.core.config import settings as _s

        result = PDFExtractor().extract_text(path)
        if result["status"] == "failed":
            raise RuntimeError(result.get("error", "PDF extraction failed"))

        text = result["text"]
        meta = result.get("metadata", {})
        ocr_pages = result.get("ocr_needed_pages", [])

        if ocr_pages:
            _pages = list(ocr_pages)
            ocr_r = run_with_retry(
                lambda: _get_ocr_service().extract_text_from_pdf_images(path, page_numbers=_pages),
                timeout=_s.OCR_TIMEOUT_SECONDS,
                max_retries=_s.PROCESSING_MAX_RETRIES,
                task_name="PDF-OCR",
            )
            if ocr_r.get("text"):
                ocr_text = ocr_r["text"]
                # For fully-scanned PDFs the digital layer is empty — use OCR
                # text as-is so the chunker receives clean structured Markdown.
                # For mixed PDFs, append OCR output after the digital content.
                if text.strip():
                    text = text + "\n\n" + ocr_text
                else:
                    text = ocr_text
                meta["ocr_pages"]      = len(ocr_pages)
                meta["ocr_confidence"] = ocr_r.get("confidence", 0)

        return _ok(text, path, **meta)

    except Exception as exc:
        logger.warning("PDFExtractor failed (%s), trying pypdf fallback", exc)
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            pages = [p.extract_text().replace("\x00", "") for p in reader.pages if p.extract_text()]
            text = "\n".join(pages).strip()
            if len(text) > 50:
                return _ok(text, path, pages=len(reader.pages), method="pypdf_fallback")
        except Exception:
            pass
        return _fail(path, str(exc))


@_register("docx", "doc")
def _extract_word(path: str) -> ExtractionResult:
    """DOCX/DOC — unstructured primary, python-docx fallback."""
    try:
        from unstructured.partition.docx import partition_docx
        elements = partition_docx(filename=path)
        parts = []
        for el in elements:
            t = str(el).strip()
            if not t:
                continue
            et = type(el).__name__
            if "Title" in et:
                parts.append(f"\n# {t}\n")
            elif "Header" in et or "Heading" in et:
                parts.append(f"\n## {t}\n")
            elif "ListItem" in et:
                parts.append(f"- {t}")
            else:
                parts.append(t)
        return _ok("\n".join(parts), path, elements=len(elements), method="unstructured")
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("unstructured DOCX failed: %s — falling back to python-docx", exc)

    from docx import Document
    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows.append("|" + "|".join(["---"] * len(cells)) + "|")
        if rows:
            parts.append("\n" + "\n".join(rows) + "\n")
    return _ok("\n".join(parts), path, paragraphs=len(doc.paragraphs), method="python-docx")


@_register("pptx", "ppt")
def _extract_pptx(path: str) -> ExtractionResult:
    """PPTX/PPT — unstructured primary, python-pptx fallback."""
    try:
        from unstructured.partition.pptx import partition_pptx
        elements = partition_pptx(filename=path)
        parts, slide_no = [], 0
        for el in elements:
            t = str(el).strip()
            if not t:
                continue
            et = type(el).__name__
            if "Title" in et:
                slide_no += 1
                parts.append(f"\n--- Slide {slide_no} ---\n")
                parts.append(f"# {t}\n")
            elif "ListItem" in et:
                parts.append(f"- {t}")
            else:
                parts.append(t)
        return _ok("\n".join(parts), path, slides=slide_no, method="unstructured")
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("unstructured PPTX failed: %s — falling back to python-pptx", exc)

    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = [f"\n--- Slide {i + 1} ---\n"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if len(slide_texts) > 1:
            parts.extend(slide_texts)
    return _ok("\n".join(parts), path, slides=len(prs.slides), method="python-pptx")


@_register("xlsx", "xls", "ods")
def _extract_spreadsheet(path: str) -> ExtractionResult:
    """Spreadsheets — structured extraction with parquet side-car."""
    import pandas as pd
    ext = Path(path).suffix.lower()
    try:
        if ext == ".ods":
            sheets = pd.read_excel(path, engine="odf", sheet_name=None)
        else:
            sheets = pd.read_excel(path, sheet_name=None)

        if not isinstance(sheets, dict):
            sheets = {"Sheet1": sheets}

        all_text_parts: list[str] = []
        parquet_paths: dict[str, str] = {}

        for sheet_name, df in sheets.items():
            # Save parquet side-car for python_tool
            parquet_path = f"{path}_{sheet_name}.parquet"
            try:
                df.to_parquet(parquet_path)
                parquet_paths[sheet_name] = parquet_path
            except Exception as exc:
                logger.warning("Parquet save failed for sheet %s: %s", sheet_name, exc)

            summary = _dataframe_summary(df, f"Sheet: {sheet_name}")
            all_text_parts.append(summary)

        total_rows = sum(len(df) for df in sheets.values())
        return _ok(
            "\n\n===\n\n".join(all_text_parts),
            path,
            rows=total_rows,
            sheets=list(sheets.keys()),
            structured_data_paths=parquet_paths,
            source_type="excel",
        )
    except Exception as exc:
        return _fail(path, f"Spreadsheet extraction failed: {exc}")


@_register("csv")
def _extract_csv(path: str) -> ExtractionResult:
    """CSV — structured extraction with parquet side-car."""
    import pandas as pd
    try:
        df = pd.read_csv(path, encoding_errors="replace")

        # Save parquet side-car for python_tool
        parquet_path = path.rsplit(".", 1)[0] + "_data.parquet"
        try:
            df.to_parquet(parquet_path)
        except Exception as exc:
            logger.warning("Parquet save failed for %s: %s", path, exc)
            parquet_path = ""

        summary = _dataframe_summary(df, Path(path).stem)
        return _ok(
            summary,
            path,
            rows=len(df),
            columns=len(df.columns),
            structured_data_path=parquet_path,
            source_type="csv",
        )
    except Exception as exc:
        return _fail(path, f"CSV extraction failed: {exc}")


@_register("txt", "md")
def _extract_text_file(path: str) -> ExtractionResult:
    """Plain text and Markdown."""
    raw = Path(path).read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    text = raw.decode(enc, errors="replace")
    return _ok(text, path, encoding=enc, method="chardet")


@_register("html", "htm")
def _extract_html(path: str) -> ExtractionResult:
    """HTML file — unstructured primary, BeautifulSoup fallback."""
    try:
        from unstructured.partition.html import partition_html
        elements = partition_html(filename=path)
        parts = []
        for el in elements:
            t = str(el).strip()
            if not t:
                continue
            et = type(el).__name__
            if "Title" in et:
                parts.append(f"\n# {t}\n")
            elif "Header" in et or "Heading" in et:
                parts.append(f"\n## {t}\n")
            elif "ListItem" in et:
                parts.append(f"- {t}")
            else:
                parts.append(t)
        if parts:
            return _ok("\n".join(parts), path, method="unstructured")
    except (ImportError, Exception):
        pass

    from bs4 import BeautifulSoup
    raw = Path(path).read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    soup = BeautifulSoup(raw.decode(enc, errors="replace"), "html.parser")
    for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
        tag.decompose()
    return _ok(soup.get_text(separator=" ", strip=True), path, method="beautifulsoup")


@_register("rtf")
def _extract_rtf(path: str) -> ExtractionResult:
    """RTF — striprtf library; falls back to raw decode."""
    try:
        from striprtf.striprtf import rtf_to_text
        raw = Path(path).read_bytes()
        enc = chardet.detect(raw).get("encoding") or "utf-8"
        rtf_string = raw.decode(enc, errors="replace")
        text = rtf_to_text(rtf_string).strip()
        if len(text) > 20:
            return _ok(text, path, method="striprtf")
    except ImportError:
        logger.debug("striprtf not installed, falling back to raw decode for RTF")
    except Exception as exc:
        logger.debug("striprtf failed: %s", exc)

    raw = Path(path).read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    text = raw.decode(enc, errors="replace")
    return _ok(text, path, method="chardet_fallback")


@_register("epub")
def _extract_epub(path: str) -> ExtractionResult:
    """EPUB — ebooklib; falls back to unzip + raw text."""
    try:
        import ebooklib
        from ebooklib import epub as _epub
        from bs4 import BeautifulSoup

        book = _epub.read_epub(path, options={"ignore_ncx": True})
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            t = soup.get_text(separator=" ", strip=True)
            if t:
                parts.append(t)
        if parts:
            title = book.get_metadata("DC", "title")
            title_str = title[0][0] if title else None
            return _ok("\n\n".join(parts), path, title=title_str, method="ebooklib")
    except ImportError:
        logger.debug("ebooklib not installed, trying zip fallback for EPUB")
    except Exception as exc:
        logger.warning("ebooklib EPUB extraction failed: %s", exc)

    # Fallback: treat EPUB as ZIP, read HTML files inside
    try:
        import zipfile
        from bs4 import BeautifulSoup
        parts = []
        with zipfile.ZipFile(path) as z:
            for name in z.namelist():
                if name.endswith((".html", ".xhtml", ".htm")):
                    raw = z.read(name)
                    soup = BeautifulSoup(raw, "html.parser")
                    t = soup.get_text(separator=" ", strip=True)
                    if t:
                        parts.append(t)
        if parts:
            return _ok("\n\n".join(parts), path, method="zip_fallback")
    except Exception as exc:
        pass

    return _fail(path, "EPUB extraction failed — install ebooklib: pip install ebooklib")


@_register("odt", "odp")
def _extract_odt(path: str) -> ExtractionResult:
    """ODT/ODP — odfpy; falls back to unzip + content.xml."""
    try:
        from odf import text as odf_text, teletype
        from odf.opendocument import load as odf_load

        doc = odf_load(path)
        paragraphs = doc.text.getElementsByType(odf_text.P)
        parts = [teletype.extractText(p) for p in paragraphs if teletype.extractText(p).strip()]
        if parts:
            return _ok("\n".join(parts), path, method="odfpy")
    except ImportError:
        logger.debug("odfpy not installed, using zip fallback for ODT")
    except Exception as exc:
        logger.warning("odfpy ODT extraction failed: %s", exc)

    # Fallback: unzip and parse content.xml
    try:
        import zipfile
        from bs4 import BeautifulSoup
        with zipfile.ZipFile(path) as z:
            if "content.xml" in z.namelist():
                xml = z.read("content.xml")
                soup = BeautifulSoup(xml, "xml")
                text = soup.get_text(separator=" ", strip=True)
                return _ok(text, path, method="zip_xml_fallback")
    except Exception:
        pass

    return _fail(path, "ODT extraction failed — install odfpy: pip install odfpy")


@_register("eml")
def _extract_eml(path: str) -> ExtractionResult:
    """EML email files — Python stdlib email module."""
    import email as _email
    from email import policy

    raw = Path(path).read_bytes()
    msg = _email.message_from_bytes(raw, policy=policy.default)

    subject = str(msg.get("Subject", ""))
    sender  = str(msg.get("From",    ""))
    date    = str(msg.get("Date",    ""))
    parts   = [f"Subject: {subject}", f"From: {sender}", f"Date: {date}", ""]

    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                try:
                    body += part.get_content() + "\n"
                except Exception:
                    pass
            elif ct == "text/html" and not body:
                from bs4 import BeautifulSoup
                try:
                    soup = BeautifulSoup(part.get_content(), "html.parser")
                    body += soup.get_text(separator=" ", strip=True) + "\n"
                except Exception:
                    pass
    else:
        try:
            body = msg.get_content()
        except Exception:
            body = raw.decode("utf-8", errors="replace")

    parts.append(body.strip())
    return _ok("\n".join(parts), path, subject=subject, method="email_stdlib")


@_register("msg")
def _extract_msg(path: str) -> ExtractionResult:
    """Outlook MSG files — extract-msg library."""
    try:
        import extract_msg
        with extract_msg.openMsg(path) as msg:
            parts = [
                f"Subject: {msg.subject or ''}",
                f"From: {msg.sender or ''}",
                f"Date: {msg.date or ''}",
                "",
                msg.body or "",
            ]
        return _ok("\n".join(parts), path, method="extract_msg")
    except ImportError:
        return _fail(path, "MSG extraction requires 'extract-msg': pip install extract-msg")
    except Exception as exc:
        return _fail(path, f"MSG extraction failed: {exc}")


# ── Generic text fallback ─────────────────────────────────────────────────────

def _generic_text_fallback(path: str) -> ExtractionResult:
    """Last-resort: try to read any file as text. Fails if alpha ratio < 5%."""
    try:
        raw = Path(path).read_bytes()
        enc = chardet.detect(raw).get("encoding") or "utf-8"
        text = raw.decode(enc, errors="replace").replace("\x00", "")

        # Quality check
        alpha = sum(c.isalpha() for c in text) / len(text) if text else 0
        if alpha < 0.05:
            return _fail(
                path,
                f"File appears to be binary (alpha ratio={alpha:.3f}). "
                "No readable text content could be extracted."
            )
        return _ok(text, path, encoding=enc, method="generic_fallback")
    except Exception as exc:
        return _fail(path, f"Generic text extraction failed: {exc}")


# ── Main extractor class ──────────────────────────────────────────────────────


class EnhancedTextExtractor:
    """Unified, format-agnostic text extractor for files, URLs, and YouTube."""

    def extract_text(self, source: str, source_type: str = "auto") -> ExtractionResult:
        """Extract text from *source* (file path or URL).

        source_type: "auto" | "file" | "url" | "youtube"
        """
        try:
            if source_type == "auto":
                source_type = self._detect_type(source)

            if source_type == "youtube":
                return self._from_youtube(source)
            if source_type == "url":
                from .youtube_service import YouTubeService
                if YouTubeService().is_youtube_url(source):
                    return self._from_youtube(source)
                return self._from_url(source)
            if source_type == "file":
                return self._from_file(source)

            raise ValueError(f"Unsupported source_type: {source_type}")

        except Exception as exc:
            logger.error("Extraction failed for %s: %s", source, exc)
            return _fail(source, str(exc), source_type)

    # ── Internal routing ──────────────────────────────────────────────────────

    @staticmethod
    def _detect_type(source: str) -> str:
        if source.startswith(("http://", "https://")):
            return "url"
        if os.path.isfile(source):
            return "file"
        raise ValueError(f"Cannot determine source type for: {source!r}")

    @staticmethod
    def _from_file(path: str) -> ExtractionResult:
        if not os.path.exists(path):
            return _fail(path, f"File not found: {path}")

        info = FileTypeDetector.detect_file_type(path)
        category = info["category"]
        ext      = info["extension"]
        mime     = info.get("mime_type", "")

        logger.info("Extracting file: %s  ext=%s  mime=%s  category=%s", path, ext, mime, category)

        # ── OCR ───────────────────────────────────────────────────────────────
        if category == "image":
            from .resilient_runner import run_with_retry
            from app.core.config import settings
            r = run_with_retry(
                lambda: _get_ocr_service().extract_text_from_image(path),
                timeout=settings.OCR_TIMEOUT_SECONDS,
                max_retries=settings.PROCESSING_MAX_RETRIES,
                task_name="OCR",
            )
            return _ok(r["text"], path, source_type="file", method="ocr",
                       confidence=r.get("confidence"))

        # ── Audio / Video → Whisper ───────────────────────────────────────────
        if category in ("audio", "video"):
            from .resilient_runner import run_with_retry
            from app.core.config import settings
            r = run_with_retry(
                lambda: _get_transcription_service().transcribe_audio_file(path),
                timeout=settings.WHISPER_TIMEOUT_SECONDS,
                max_retries=settings.PROCESSING_MAX_RETRIES,
                task_name="Whisper",
            )
            return _ok(r["text"], path, source_type="file", method="transcription",
                       language=r.get("language"), duration=r.get("duration"))

        # ── Document — registered extractor ──────────────────────────────────
        extractor_fn = _DOC_EXTRACTORS.get(ext)
        if extractor_fn:
            try:
                return extractor_fn(path)
            except Exception as exc:
                logger.warning("Registered extractor for %s failed: %s — trying generic fallback", ext, exc)
                return _generic_text_fallback(path)

        # ── Unknown type — generic text fallback ──────────────────────────────
        logger.info("No registered extractor for ext=%s — trying generic text fallback", ext)
        return _generic_text_fallback(path)

    @staticmethod
    def _from_url(url: str) -> ExtractionResult:
        from .web_scraping import WebScrapingService
        scraper = WebScrapingService()

        # 1. Detect URL content type via HTTP headers
        url_info  = scraper.detect_url_type(url)
        category  = url_info.get("category", "unknown")

        # 2. If header detection failed / returned unknown, try URL extension
        if category == "unknown" or category == "web":
            ext_cat = FileTypeDetector.detect_from_extension(url)
            if ext_cat and ext_cat not in ("web",):
                logger.info(
                    "Header detection gave category=%s; URL extension suggests %s — trying download",
                    category, ext_cat,
                )
                category = ext_cat

        # 3. If it's a downloadable file type, download → extract as file
        exts = FileTypeDetector.get_supported_extensions()
        _DOWNLOAD_CATS = set(exts) | {"image", "audio", "video", "document"}
        
        if category in _DOWNLOAD_CATS:
            logger.info("URL category=%s — downloading for direct extraction: %s", category, url)
            temp_path = scraper.download_url_to_temp(url)
            if temp_path:
                try:
                    result = EnhancedTextExtractor._from_file(temp_path)
                    result["source"]      = url
                    result["source_type"] = "url"
                    return result
                except Exception as exc:
                    logger.warning("Direct file extraction failed for URL %s: %s", url, exc)
                finally:
                    try:
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                    except OSError:
                        pass

        # 4. Fall back to web scraping
        r = scraper.extract_content_from_url(url)
        if r["status"] == "failed":
            return _fail(url, r.get("error", "Web scraping failed"), "url")
        return _ok(r["text"], url, "url", title=r.get("title"))

    @staticmethod
    def _from_youtube(url: str) -> ExtractionResult:
        from .youtube_service import YouTubeService
        r = YouTubeService().extract_transcript_from_url(url)
        text = r.get("transcript", "")
        if not text and r.get("description"):
            text = f"[Video Title: {r.get('title', '')}]\n\n{r['description']}"
        if not text:
            return _fail(url, r.get("error", "No transcript available"), "youtube")
        return _ok(
            text, url, "youtube",
            title=r.get("title"),
            duration=r.get("duration"),
            transcript_lang=r.get("transcript_language"),
        )


# ── Legacy shim ───────────────────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    """Legacy helper — returns plain text string."""
    return EnhancedTextExtractor().extract_text(file_path).get("text", "")
